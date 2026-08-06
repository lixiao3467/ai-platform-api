"""Document parsers — extract text from various file formats."""

from __future__ import annotations

import csv
import io
import base64

import structlog

logger = structlog.get_logger()


# MIME type → parser dispatch table
_MIME_DISPATCH: dict[str, str] = {
    # Text formats (handled inline)
    "text/plain": "_text",
    "text/markdown": "_text",
    "": "_text",
    # PDF
    "application/pdf": "_parse_pdf",
    # HTML
    "text/html": "_parse_html",
    # Word
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "_parse_docx",
    # Excel
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "_parse_excel",
    "application/vnd.ms-excel": "_parse_excel",
    # CSV
    "text/csv": "_parse_csv",
    # PPT
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "_parse_pptx",
    # Images
    "image/png": "_parse_image",
    "image/jpeg": "_parse_image",
    "image/jpg": "_parse_image",
    "image/webp": "_parse_image",
}


async def parse_document(content: bytes, mime_type: str) -> str:
    """
    Parse a document into plain text based on MIME type.

    Supported formats:
    - text/plain, text/markdown
    - application/pdf
    - text/html
    - Word (.docx)
    - Excel (.xlsx, .xls legacy handled as xlsx best-effort)
    - CSV
    - PPT (.pptx)
    - Images (png/jpeg/jpg/webp) via multimodal LLM OCR
    """
    parser_name = _MIME_DISPATCH.get(mime_type)

    if parser_name is None:
        # Fallback: try as plain text
        return content.decode("utf-8", errors="replace")

    if parser_name == "_text":
        return content.decode("utf-8", errors="replace")

    parser = globals()[parser_name]
    try:
        result = await parser(content)
        return result
    except Exception as e:
        logger.warning(
            "Parser failed",
            parser=parser_name,
            mime_type=mime_type,
            error=str(e),
        )
        return f"[文档解析失败: {parser_name} — {type(e).__name__}: {e}]"


async def _parse_pdf(content: bytes) -> str:
    """Parse PDF using pymupdf4llm."""
    try:
        import pymupdf4llm

        text = pymupdf4llm.to_markdown(content)
        return text
    except ImportError:
        # Fallback to PyMuPDF directly
        import fitz

        doc = fitz.open(stream=content, filetype="pdf")
        text = "\n\n".join(page.get_text() for page in doc)
        doc.close()
        return text


async def _parse_html(content: bytes) -> str:
    """Parse HTML using BeautifulSoup."""
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(content, "html.parser")
    # Remove script and style elements
    for element in soup(["script", "style"]):
        element.decompose()
    return soup.get_text(separator="\n", strip=True)


async def _parse_docx(content: bytes) -> str:
    """Parse DOCX using python-docx."""
    from docx import Document

    doc = Document(io.BytesIO(content))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


# =============================================================================
# Excel parser
# =============================================================================


async def _parse_excel(content: bytes) -> str:
    """Parse Excel (.xlsx/.xls) using openpyxl. Each sheet → markdown table."""
    from openpyxl import load_workbook

    try:
        wb = load_workbook(io.BytesIO(content), data_only=True)
    except Exception as e:
        return f"[Excel 文件解析失败: {e}]"

    max_rows = 1000
    output: list[str] = []

    try:
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]

            # Unfill merged cells: propagate top-left value across the merge range
            # so we render a flat rectangular grid.
            merged_ranges = list(ws.merged_cells.ranges)

            rows: list[list[str]] = []
            truncated = False

            for row_idx, row in enumerate(ws.iter_rows(values_only=False), start=1):
                if row_idx > max_rows + 1:
                    truncated = True
                    break
                cells: list[str] = []
                for cell in row:
                    val = cell.value
                    # Resolve merged cells: if cell is inside a merged range
                    # but isn't the top-left, use the top-left's value.
                    if val is None:
                        for mr in merged_ranges:
                            if cell.coordinate in mr and cell.coordinate != mr.start_cell.coordinate:
                                val = ws[mr.start_cell.coordinate].value
                                break
                    cells.append("" if val is None else str(val))
                rows.append(cells)

            if not rows:
                output.append(f"## Sheet: {sheet_name}\n\n_(空表)_\n")
                continue

            # Determine column count (max width across rows)
            col_count = max(len(r) for r in rows) if rows else 0
            # Pad rows to uniform width
            for r in rows:
                if len(r) < col_count:
                    r.extend([""] * (col_count - len(r)))

            # Build markdown table
            md = [f"## Sheet: {sheet_name}", ""]
            # Header = first row
            header = rows[0]
            md.append("| " + " | ".join(_md_escape(c) for c in header) + " |")
            md.append("|" + "|".join(["---------"] * col_count) + "|")
            for r in rows[1:]:
                md.append("| " + " | ".join(_md_escape(c) for c in r) + " |")

            if truncated:
                md.append(f"\n_(表格已截断，仅显示前 {max_rows} 行，共 {ws.max_row or '?'} 行)_")

            output.append("\n".join(md))
    finally:
        wb.close()

    return "\n\n".join(output).strip()


def _md_escape(value: str) -> str:
    """Escape pipe characters and newlines for markdown table cells."""
    return value.replace("|", "\\|").replace("\n", " ").replace("\r", "")


# =============================================================================
# CSV parser
# =============================================================================


_CSV_DELIMITERS = [",", "\t", ";"]


async def _parse_csv(content: bytes) -> str:
    """Parse CSV. Auto-detect delimiter. Output markdown table."""
    text = _decode_text(content)

    # Detect delimiter via csv.Sniffer on a sample
    sample = text[:8192]
    delimiter = ","
    try:
        dialect = csv.Sniffer().sniff(sample, delimiters="".join(_CSV_DELIMITERS))
        delimiter = dialect.delimiter
    except csv.Error:
        # Fallback: count occurrences
        counts = {d: sample.count(d) for d in _CSV_DELIMITERS}
        delimiter = max(counts, key=lambda d: counts[d])

    reader = csv.reader(io.StringIO(text), delimiter=delimiter)
    max_rows = 1000
    rows: list[list[str]] = []
    truncated = False
    total_rows = 0

    for row in reader:
        total_rows += 1
        if len(rows) >= max_rows:
            truncated = True
            continue
        rows.append(row)

    if not rows:
        return "_(CSV 文件为空)_"

    col_count = max(len(r) for r in rows)
    for r in rows:
        if len(r) < col_count:
            r.extend([""] * (col_count - len(r)))

    md: list[str] = []
    header = rows[0]
    md.append("| " + " | ".join(_md_escape(c) for c in header) + " |")
    md.append("|" + "|".join(["---------"] * col_count) + "|")
    for r in rows[1:]:
        md.append("| " + " | ".join(_md_escape(c) for c in r) + " |")

    if truncated:
        md.append(f"\n_(表格已截断，仅显示前 {max_rows} 行，共 {total_rows} 行)_")

    return "\n".join(md)


def _decode_text(content: bytes) -> str:
    """Decode bytes with fallback chain: UTF-8 → GBK → Latin-1."""
    for enc in ("utf-8", "gbk", "gb2312", "latin-1"):
        try:
            return content.decode(enc)
        except UnicodeDecodeError:
            continue
    return content.decode("utf-8", errors="replace")


# =============================================================================
# PPT parser
# =============================================================================


async def _parse_pptx(content: bytes) -> str:
    """Parse PPTX using python-pptx. Each slide → heading + body."""
    from pptx import Presentation
    from pptx.util import Inches  # noqa: F401  (kept for future use)

    try:
        prs = Presentation(io.BytesIO(content))
    except Exception as e:
        return f"[PPT 文件解析失败: {e}]"

    output: list[str] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        # Determine title
        title_text = ""
        if slide.shapes.title is not None:
            title_text = slide.shapes.title.text.strip()

        heading = f"## Slide {slide_idx}: {title_text}" if title_text else f"## Slide {slide_idx}"

        body_parts: list[str] = []

        for shape in slide.shapes:
            # Text frames (skip the title shape — already captured above)
            if shape.has_text_frame and shape != slide.shapes.title:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        body_parts.append(text)

            # Tables
            if shape.has_table:
                table = shape.table
                rows_data: list[list[str]] = []
                for row in table.rows:
                    rows_data.append([cell.text.strip() for cell in row.cells])
                if rows_data:
                    col_count = max(len(r) for r in rows_data)
                    for r in rows_data:
                        while len(r) < col_count:
                            r.append("")
                    header = rows_data[0]
                    tbl_lines = [
                        "| " + " | ".join(_md_escape(c) for c in header) + " |",
                        "|" + "|".join(["---------"] * col_count) + "|",
                    ]
                    for r in rows_data[1:]:
                        tbl_lines.append("| " + " | ".join(_md_escape(c) for c in r) + " |")
                    body_parts.append("\n".join(tbl_lines))

            # Images with alt text
            if shape.shape_type is not None and hasattr(shape, "image"):
                try:
                    alt = shape.name or ""
                    if alt:
                        body_parts.append(f"_[图片: {alt}]_")
                except Exception:
                    pass

        slide_text = heading
        if body_parts:
            slide_text += "\n\n" + "\n\n".join(body_parts)
        output.append(slide_text)

    return "\n\n".join(output).strip() if output else "_(PPT 文件为空)_"


# =============================================================================
# Image OCR parser (via multimodal LLM)
# =============================================================================


async def _parse_image(content: bytes) -> str:
    """
    Parse image via multimodal LLM (vision OCR).

    Uses the already-configured LiteLLM gateway. Falls back gracefully
    if the LLM call fails — returns a placeholder string rather than
    raising, so the engine can record the state.
    """
    try:
        import litellm

        from ai_platform.config import get_settings

        settings = get_settings()
        model = settings.vision_model or "gpt-4o"

        b64 = base64.b64encode(content).decode("ascii")
        # Pick a safe MIME type — default to jpeg; png/webp also work with
        # most vision models.
        data_url = f"data:image/jpeg;base64,{b64}"

        response = await litellm.acompletion(
            model=model,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": (
                                "请提取这张图片中的所有文字内容。如果是图表，请描述图表内容。"
                                "直接输出识别到的文本，不要添加额外解释。"
                            ),
                        },
                        {
                            "type": "image_url",
                            "image_url": {"url": data_url},
                        },
                    ],
                }
            ],
        )

        text = response["choices"][0]["message"]["content"]
        if not text or not text.strip():
            return "[图片内容无法提取]"
        return text.strip()

    except Exception as e:
        logger.warning("Image OCR failed", error=str(e))
        return "[图片内容无法提取]"
